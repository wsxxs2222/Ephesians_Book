import json
from pptx import Presentation
import util

processed_data_folder_path = './processed_data/'
raw_data_folder_path = "./raw_data/"

def file_name_builder(folder_path, file_prefix, file_extension, index):
  return f"{folder_path}{file_prefix}_{index:02d}.{file_extension}"

def convert_ppt_files_to_title_json(start_index, end_index):
    for i in range(start_index, end_index):
        ppt_file = util.file_name_builder(f"{raw_data_folder_path}ppt/", "Eph"
        , "pptx", i)
        prs = Presentation(ppt_file)
        
        title_object_list = []
        
        for slide_idx, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        title_object_list.append({
                            "content": text,
                            "level": para.level,
                            "contentType": "",   # placeholder for manual annotation
                            "slideNumber": slide_idx
                        })
        
        json_file_name = util.file_name_builder(f"{processed_data_folder_path}title_from_ppt/", "Eph"
        , "json", i)
        # Save as JSON file
        with open(json_file_name, "w", encoding="utf-8") as f:
            json.dump(title_object_list, f, ensure_ascii=False, indent=2)